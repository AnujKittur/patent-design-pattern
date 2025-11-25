"""
Generate PowerPoint presentation for Patent Design Pattern Project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Create comprehensive PowerPoint presentation."""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    title_color = RGBColor(25, 118, 210)  # Blue
    accent_color = RGBColor(32, 56, 88)   # Dark blue
    text_color = RGBColor(51, 51, 51)     # Dark gray
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background rectangle
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Inches(2)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.color.rgb = accent_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Patent Design Pattern Project"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "RAG-Based Construction Robotics Design Generation System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = text_color
    subtitle_para.alignment = PP_ALIGN.LEFT
    
    # Author/Date
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    author_frame = author_box.text_frame
    author_frame.text = "End-to-End Project Presentation"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(18)
    author_para.font.color.rgb = text_color
    author_para.alignment = PP_ALIGN.LEFT
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Project Overview"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "RAG-Based System for Construction Robotics Design"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Retrieval-Augmented Generation (RAG) architecture"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Patent-grounded design generation"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• ~200 construction robotics patents from USPTO PatentsView API"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Hybrid retrieval: BM25 (keyword) + Vector (semantic) search"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Citation-grounded: Every design choice must cite patent sources"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Structured JSON output: Modules, Actuation, Sensing, Control, Materials, Safety, BOM"
    p.font.size = Pt(16)
    p.level = 0
    
    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Statement"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Challenges in Construction Robotics Design"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Vast repositories of patent data exist but are difficult to process"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Manual extraction of design patterns is time-consuming and error-prone"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Need for automated systems that can leverage patent knowledge"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Design choices should be grounded in existing patent literature"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Requirements: Fast, accurate, and traceable design generation"
    p.font.size = Pt(16)
    p.level = 0
    
    # Slide 4: Solution Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Solution Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "RAG Pipeline: Retrieve → Rerank → Generate"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "1. Ingestion Layer"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • USPTO PatentsView API → Download patents"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Chunk documents (abstract, description, claims)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Generate embeddings → ChromaDB vector index"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Retrieval Layer"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Multi-query expansion (GPT-4o-mini)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Hybrid search: 40% BM25 + 60% Vector similarity"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Cross-encoder reranking (BAAI/bge-reranker-large)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Generation Layer"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • GPT-4o-mini with citation-enforced prompts"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Structured JSON output with citations"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 5: System Components
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Components"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Core Implementation Files"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "ingest.py - Patent ingestion and indexing"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Downloads patents from PatentsView API"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Chunks documents (abstract, description, claims)"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Generates embeddings and builds ChromaDB index"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "app.py - FastAPI server"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • POST /design endpoint for design generation"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Hybrid retrieval with BM25 + Vector search"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • LLM-based generation with citation enforcement"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "streamlit_app.py - Web UI"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Interactive design generation interface"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Visual design brief display with citations"
    p.font.size = Pt(12)
    p.level = 1
    
    # Slide 6: Data Flow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Flow Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "End-to-End Pipeline"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "USPTO PatentsView API"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "    ↓"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Patent Download & Processing (ingest.py)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "    ↓"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Chunking & Embedding Generation"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "    ↓"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "ChromaDB Vector Index + BM25 Index"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "    ↓"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "User Query → Multi-Query Expansion"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "    ↓"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Hybrid Retrieval (BM25 + Vector)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "    ↓"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Reranking → Top 10 Patents"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "    ↓"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "LLM Generation → Structured Design Brief"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    # Slide 7: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "System Capabilities"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "✓ Citation-Grounded Design"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   Every design choice must cite at least one patent source"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Hybrid Retrieval Strategy"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   Combines keyword (BM25) and semantic (vector) search for better results"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Multi-Query Expansion"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   Generates diverse query variations to improve retrieval coverage"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Cross-Encoder Reranking"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   Fine-tunes retrieval results using relevance scoring"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Structured Output"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   JSON format with modules, actuation, sensing, control, materials, safety, BOM"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 8: Technical Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technical Stack"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Technologies & Libraries"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Backend Framework: FastAPI + Uvicorn"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Vector Database: ChromaDB (persistent storage)"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Embeddings: OpenAI text-embedding-3-large (or sentence-transformers)"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "LLM: OpenAI GPT-4o-mini (or local fallback)"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Keyword Search: rank-bm25"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Reranking: sentence-transformers CrossEncoder (BAAI/bge-reranker-large)"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Web UI: Streamlit"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Data Processing: requests, BeautifulSoup, pypdfium2, PIL"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Evaluation: RAGAS framework"
    p.font.size = Pt(14)
    p.level = 0
    
    # Slide 9: Patent Data Structure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Patent Data Structure"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Data Model & Chunking Strategy"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Patent Document Fields:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • patent_number, title, abstract"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • claims_text, description"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • cpc (classification codes), assignee, pub_year, url"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Chunking Strategy:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   1. Abstract: Single chunk"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   2. Description: 800-1200 tokens with 120 token overlap"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   3. Claims: One chunk per claim"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Metadata Enrichment:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • CPC codes, publication year, mechanism tags"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Figure paths for visualization"
    p.font.size = Pt(12)
    p.level = 1
    
    # Slide 10: Retrieval Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Retrieval Strategy"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Hybrid Search Approach"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Step 1: Multi-Query Expansion"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Generate 3 diverse query variations using GPT-4o-mini"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Focus on different aspects of the design specification"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Step 2: BM25 Keyword Search"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Tokenize queries and corpus"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Score documents based on term frequency"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Step 3: Vector Semantic Search"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Average embeddings from expanded queries"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Cosine similarity search in ChromaDB"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Step 4: Score Fusion"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Normalize both scores to [0, 1]"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Weighted combination: 40% BM25 + 60% Vector"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Step 5: Reranking"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Cross-encoder reranking with BAAI/bge-reranker-large"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Return top 10 most relevant patents"
    p.font.size = Pt(12)
    p.level = 1
    
    # Slide 11: Design Generation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Design Generation"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "LLM-Based Design Synthesis"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Prompt Engineering:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • System prompt: Enforces citation requirements"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Designer prompt: Structured template with retrieved context"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Context formatting: [US<patent>, section] + patent text"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Output Structure:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Overview: High-level system description"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Modules: System components with functions and citations"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Actuation, Sensing, Control: Technical subsystems"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Materials, Safety, Procedure: Design specifications"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • BOM: Bill of materials with quantities and costs"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Citations: Full patent references with URLs"
    p.font.size = Pt(12)
    p.level = 1
    
    # Slide 12: User Interfaces
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "User Interfaces"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Multiple Access Methods"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "1. Streamlit Web UI (streamlit_app.py)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Interactive design specification input"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Optional filters: CPC codes, year range"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Visual design brief display with expandable sections"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Patent citations with clickable links"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Optional image generation for design visualization"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. FastAPI REST API (app.py)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • POST /design endpoint"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • JSON request/response format"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Swagger/OpenAPI documentation at /docs"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. CLI Tool (demo.py)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Command-line interface for batch processing"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • JSON or pretty-print output"
    p.font.size = Pt(12)
    p.level = 1
    
    # Slide 13: Citation System
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Citation System"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Ensuring Traceability"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Citation Format:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Standard: (US<patent_number>, section)"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • With claim: (US<patent_number>, section, claim_1)"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Speculative: SPECULATIVE (if no patent evidence)"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Citation Requirements:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Every design choice must have at least one citation"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Modules, actuation, sensing, control, materials, safety, BOM items"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Full citation list with patent numbers, titles, URLs, relevance"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Enforcement:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • System prompt explicitly requires citations"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Designer prompt includes citation examples"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Context format makes citations easy to extract"
    p.font.size = Pt(12)
    p.level = 1
    
    # Slide 14: Evaluation Framework
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Evaluation Framework"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "RAGAS Metrics"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Faithfulness: Measures factual consistency"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Ensures generated design is grounded in retrieved patents"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Answer Relevancy: Measures relevance to query"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Evaluates if design meets user specifications"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Context Precision: Measures precision of retrieved context"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Assesses quality of retrieved patent documents"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Context Recall: Measures recall of retrieved context"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Ensures all relevant patents are retrieved"
    p.font.size = Pt(12)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Evaluation Script: eval_ragas.py"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    # Slide 15: Project Structure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Structure"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "File Organization"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Core Scripts:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • ingest.py - Patent ingestion and indexing"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • app.py - FastAPI server"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • demo.py - CLI interface"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • streamlit_app.py - Web UI"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • eval_ragas.py - Evaluation script"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Directories:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • data/raw/ - Raw patent JSONL files"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • data/chunks/ - Chunked patent data"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • data/index/ - ChromaDB vector index"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • data/media/ - Patent figures and images"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Configuration:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • prompts/ - System and designer prompt templates"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • requirements.txt - Python dependencies"
    p.font.size = Pt(11)
    p.level = 1
    
    # Slide 16: Usage Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Usage Workflow"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Getting Started"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "1. Install Dependencies"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   pip install -r requirements.txt"
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Set Up Environment"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   Create .env file with OPENAI_API_KEY (optional)"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Ingest Patents"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   python ingest.py"
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   Downloads ~200 patents, chunks, and builds index"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Start API Server"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   python app.py"
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5. Use Interface (Choose One)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Streamlit: streamlit run streamlit_app.py"
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • CLI: python demo.py --prompt \"...\""
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.level = 1
    
    # Slide 17: Example Output
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Example Output"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Design Brief Structure"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Input: \"Design a robotic bricklaying system for 3-story buildings\""
    p.font.size = Pt(12)
    p.font.italic = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(8)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Output Sections:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Overview: High-level system description"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Modules: [Robotic Arm, Control System, Material Handling, ...]"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Actuation: [Arm Actuation: Electric servo motors, ...]"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Sensing: [Position Sensing: Encoders + Camera, ...]"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Control: [Motion Control: PID, Autonomous, ...]"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Materials: [Frame: Aluminum, ...]"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Safety: [Collision: Emergency stop + barriers, ...]"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Procedure: [Step-by-step implementation]"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • BOM: [Items with quantities, costs, citations]"
    p.font.size = Pt(11)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Citations: [Full patent references with URLs]"
    p.font.size = Pt(11)
    p.level = 1
    
    # Slide 18: Key Achievements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Achievements"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Project Highlights"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "✓ Complete End-to-End RAG System"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   From patent ingestion to design generation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Hybrid Retrieval Implementation"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   Combines BM25 and vector search for optimal results"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Citation-Grounded Output"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   Every design choice traceable to patent sources"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Multiple Interface Options"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   Web UI, REST API, and CLI for different use cases"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Production-Ready Architecture"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   Scalable, modular, and well-documented"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 19: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Enhancements"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Potential Improvements"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Scale to larger patent database (1000+ patents)"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Enhanced reranking with cross-encoders"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Multi-modal support (patent figures, diagrams)"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Fine-tuned LLM models on patent domain"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Distributed ChromaDB for better scalability"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Multilingual support (patent documents in different languages)"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Advanced evaluation metrics and benchmarking"
    p.font.size = Pt(14)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Real-time patent updates and incremental indexing"
    p.font.size = Pt(14)
    p.level = 0
    
    # Slide 20: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Summary"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "A complete RAG-based system for generating construction robotics designs"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "grounded in patent literature, with hybrid retrieval, citation enforcement,"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "and multiple interface options for seamless integration."
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Key Value Propositions:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "   • Automates design pattern extraction from patents"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Ensures design choices are grounded in existing literature"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Provides structured, comprehensive design briefs"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Supports multiple access methods (Web, API, CLI)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Thank You!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.level = 0
    
    # Save presentation
    output_path = "Patent_Design_Pattern_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation created successfully: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_presentation()

